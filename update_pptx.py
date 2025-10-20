from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

pptx_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Portfolio.pptx"

try:
    prs = Presentation(pptx_path)

    # Get the layout from an existing slide (using slide 3's layout as reference)
    slide_layout = prs.slides[2].slide_layout

    # Insert new slide at position 3 (which will be slide 4, after current slide 3)
    new_slide = prs.slides.add_slide(slide_layout)

    # Clear existing shapes from the layout
    for shape in new_slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()

    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(0.8)

    title_box = new_slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "Ainex - International Expansion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Add subtitle
    subtitle_left = Inches(0.5)
    subtitle_top = Inches(1.4)
    subtitle_width = Inches(9)
    subtitle_height = Inches(0.6)

    subtitle_box = new_slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "ENAD System Deployment & Technical Support"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)

    # Add main content
    content_left = Inches(0.8)
    content_top = Inches(2.2)
    content_width = Inches(8.4)
    content_height = Inches(4.5)

    content_box = new_slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Project Overview
    p1 = content_frame.paragraphs[0]
    p1.text = "Project Overview"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 51, 102)
    p1.space_after = Pt(10)

    # Description
    p2 = content_frame.add_paragraph()
    p2.text = "Led international expansion of ENAD AI-powered endoscopy diagnostic system to Southeast Asian markets"
    p2.font.size = Pt(14)
    p2.space_after = Pt(18)
    p2.level = 0

    # Target Markets header
    p3 = content_frame.add_paragraph()
    p3.text = "Target Markets & Activities"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0, 51, 102)
    p3.space_after = Pt(10)

    # Thailand
    p4 = content_frame.add_paragraph()
    p4.text = "Thailand"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.space_after = Pt(6)
    p4.level = 0

    p5 = content_frame.add_paragraph()
    p5.text = "Hospital product demonstrations and technical support"
    p5.font.size = Pt(13)
    p5.space_after = Pt(6)
    p5.level = 1

    p6 = content_frame.add_paragraph()
    p6.text = "ENAD system integration and clinician training"
    p6.font.size = Pt(13)
    p6.space_after = Pt(12)
    p6.level = 1

    # Singapore
    p7 = content_frame.add_paragraph()
    p7.text = "Singapore"
    p7.font.size = Pt(14)
    p7.font.bold = True
    p7.space_after = Pt(6)
    p7.level = 0

    p8 = content_frame.add_paragraph()
    p8.text = "Advanced medical facility system deployment"
    p8.font.size = Pt(13)
    p8.space_after = Pt(6)
    p8.level = 1

    p9 = content_frame.add_paragraph()
    p9.text = "Technical consultation and ongoing support"
    p9.font.size = Pt(13)
    p9.space_after = Pt(12)
    p9.level = 1

    # Vietnam
    p10 = content_frame.add_paragraph()
    p10.text = "Vietnam"
    p10.font.size = Pt(14)
    p10.font.bold = True
    p10.space_after = Pt(6)
    p10.level = 0

    p11 = content_frame.add_paragraph()
    p11.text = "Regional hospital network product introduction"
    p11.font.size = Pt(13)
    p11.space_after = Pt(6)
    p11.level = 1

    p12 = content_frame.add_paragraph()
    p12.text = "Post-deployment technical assistance and maintenance"
    p12.font.size = Pt(13)
    p12.level = 1

    # Move the new slide to position 3 (after current slide 3, before current slide 4)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])  # Remove the last slide (newly added)
    xml_slides.insert(3, slides[-1])  # Insert it at position 3

    # Save the presentation
    output_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Portfolio.pptx"
    prs.save(output_path)

    print(f"[SUCCESS] Successfully added new slide at position 4")
    print(f"[SUCCESS] Slide title: 'Ainex - International Expansion'")
    print(f"[SUCCESS] Content: ENAD system deployment in Thailand, Singapore, and Vietnam")
    print(f"[SUCCESS] Saved to: {output_path}")

except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
