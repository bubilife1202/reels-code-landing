from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

pptx_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Portfolio.pptx"

try:
    prs = Presentation(pptx_path)

    # Analyze slide 3 (Ainex Corporation slide) as reference
    slide_3 = prs.slides[2]

    print("=== SLIDE 3 (Ainex Corporation) ANALYSIS ===\n")

    for idx, shape in enumerate(slide_3.shapes):
        print(f"Shape {idx}:")
        print(f"  Type: {shape.shape_type}")
        print(f"  Has text: {hasattr(shape, 'text')}")

        if hasattr(shape, "text") and shape.text.strip():
            print(f"  Text: {shape.text[:50]}...")

        if hasattr(shape, "text_frame"):
            tf = shape.text_frame
            print(f"  Position: Left={shape.left}, Top={shape.top}")
            print(f"  Size: Width={shape.width}, Height={shape.height}")

            for p_idx, para in enumerate(tf.paragraphs):
                if para.text.strip():
                    print(f"  Paragraph {p_idx}:")
                    print(f"    Text: {para.text[:50]}")
                    print(f"    Level: {para.level}")
                    if para.runs:
                        run = para.runs[0]
                        print(f"    Font size: {run.font.size}")
                        print(f"    Font bold: {run.font.bold}")
                        if run.font.color.rgb:
                            print(f"    Font color: {run.font.color.rgb}")
        print()

except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
