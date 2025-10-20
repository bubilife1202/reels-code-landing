from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

pptx_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Portfolio.pptx"

try:
    prs = Presentation(pptx_path)

    # Delete the slide we added earlier (slide 3, index 3)
    rId = prs.slides._sldIdLst[3].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[3]

    # Get the layout from slide 3 (Ainex Corporation)
    slide_layout = prs.slides[2].slide_layout

    # Add new slide
    new_slide = prs.slides.add_slide(slide_layout)

    # Clear existing shapes
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Add background shape (dark background at top)
    bg_shape = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        Emu(9144000), Emu(1219200)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(28, 40, 51)  # Dark blue-gray
    bg_shape.line.fill.background()

    # Add main title "Ainex Corporation"
    title_shape = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(304800), Emu(228600),
        Emu(4362260), Emu(381000)
    )
    title_shape.fill.background()
    title_shape.line.fill.background()
    tf = title_shape.text_frame
    tf.text = "Ainex Corporation"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add subtitle "International Expansion"
    subtitle_shape = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(304800), Emu(685800),
        Emu(8705088), Emu(266700)
    )
    subtitle_shape.fill.background()
    subtitle_shape.line.fill.background()
    tf = subtitle_shape.text_frame
    tf.text = "International Expansion - ENAD System Deployment"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(170, 183, 184)

    # Add "Key Activities" header
    header_shape = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(457200), Emu(1504950),
        Emu(8394192), Emu(304800)
    )
    header_shape.fill.background()
    header_shape.line.fill.background()
    tf = header_shape.text_frame
    tf.text = "Key Activities"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(93, 173, 226)

    # Add description
    desc_shape = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(457200), Emu(1962150),
        Emu(8394192), Emu(609600)
    )
    desc_shape.fill.background()
    desc_shape.line.fill.background()
    tf = desc_shape.text_frame
    tf.text = "Led international expansion of ENAD AI-powered endoscopy diagnostic system to Southeast Asian healthcare markets"
    p = tf.paragraphs[0]
    p.font.size = Pt(13.5)
    p.font.color.rgb = RGBColor(46, 64, 83)

    # Left box - Thailand
    left_box = new_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(457200), Emu(2838450),
        Emu(3962400), Emu(2019300)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(245, 247, 249)
    left_box.line.color.rgb = RGBColor(93, 173, 226)
    left_box.line.width = Pt(2)

    # Thailand title
    thai_title = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(685800), Emu(3067050),
        Emu(3575304), Emu(266700)
    )
    thai_title.fill.background()
    thai_title.line.fill.background()
    tf = thai_title.text_frame
    tf.text = "Thailand"
    p = tf.paragraphs[0]
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(28, 40, 51)

    # Thailand content
    thai_content = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(685800), Emu(3486150),
        Emu(3505200), Emu(1143000)
    )
    thai_content.fill.background()
    thai_content.line.fill.background()
    tf = thai_content.text_frame

    p1 = tf.paragraphs[0]
    p1.text = "Hospital product demonstrations"
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(46, 64, 83)

    p2 = tf.add_paragraph()
    p2.text = "ENAD system integration and training"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(46, 64, 83)

    p3 = tf.add_paragraph()
    p3.text = "On-site technical support"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(46, 64, 83)

    # Right box - Singapore
    right_box = new_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(4724400), Emu(2838450),
        Emu(3962400), Emu(950000)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(245, 247, 249)
    right_box.line.color.rgb = RGBColor(93, 173, 226)
    right_box.line.width = Pt(2)

    # Singapore title
    sg_title = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(4953000), Emu(3067050),
        Emu(3575304), Emu(266700)
    )
    sg_title.fill.background()
    sg_title.line.fill.background()
    tf = sg_title.text_frame
    tf.text = "Singapore"
    p = tf.paragraphs[0]
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(28, 40, 51)

    # Singapore content
    sg_content = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(4953000), Emu(3486150),
        Emu(3505200), Emu(876300)
    )
    sg_content.fill.background()
    sg_content.line.fill.background()
    tf = sg_content.text_frame

    p1 = tf.paragraphs[0]
    p1.text = "Advanced facility system deployment"
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(46, 64, 83)

    p2 = tf.add_paragraph()
    p2.text = "Technical consultation and support"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(46, 64, 83)

    # Bottom right box - Vietnam
    vietnam_box = new_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(4724400), Emu(3971800),
        Emu(3962400), Emu(886000)
    )
    vietnam_box.fill.solid()
    vietnam_box.fill.fore_color.rgb = RGBColor(245, 247, 249)
    vietnam_box.line.color.rgb = RGBColor(93, 173, 226)
    vietnam_box.line.width = Pt(2)

    # Vietnam title
    vn_title = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(4953000), Emu(4200400),
        Emu(3575304), Emu(266700)
    )
    vn_title.fill.background()
    vn_title.line.fill.background()
    tf = vn_title.text_frame
    tf.text = "Vietnam"
    p = tf.paragraphs[0]
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(28, 40, 51)

    # Vietnam content
    vn_content = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(4953000), Emu(4619500),
        Emu(3505200), Emu(810300)
    )
    vn_content.fill.background()
    vn_content.line.fill.background()
    tf = vn_content.text_frame

    p1 = tf.paragraphs[0]
    p1.text = "Regional hospital network introduction"
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(46, 64, 83)

    p2 = tf.add_paragraph()
    p2.text = "Post-deployment technical assistance"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(46, 64, 83)

    # Move the new slide to position 3 (after current slide 3, before current slide 4)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(3, slides[-1])

    # Save
    output_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Portfolio.pptx"
    prs.save(output_path)

    print("[SUCCESS] Updated slide with matching template style")
    print("[SUCCESS] Slide added at position 4")
    print("[SUCCESS] Saved to:", output_path)

except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
