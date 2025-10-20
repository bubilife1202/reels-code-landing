from pptx import Presentation
import sys

pptx_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Portfolio.pptx"

try:
    prs = Presentation(pptx_path)

    print(f"Total slides: {len(prs.slides)}\n")
    print("="*80)

    for idx, slide in enumerate(prs.slides, 1):
        print(f"\n[SLIDE {idx}]")
        print("-"*80)

        # Extract all text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(text)

        print("-"*80)

except Exception as e:
    print(f"Error: {e}")
    sys.exit(1)
