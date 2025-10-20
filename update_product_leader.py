from pptx import Presentation
from pptx.util import Pt
import sys
import io

# Set UTF-8 encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

pptx_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Product_Leader.pptx"

try:
    prs = Presentation(pptx_path)

    # ============================================
    # Update Slide 6 (Phase 3: Validation) - Index 5
    # ============================================
    slide_6 = prs.slides[5]

    print("Updating Slide 6 (Phase 3: Validation)...")

    # Find and update the regulatory approval text
    for shape in slide_6.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                # Update Regulatory Approval section
                if "HSA 승인 (싱가포르)" in paragraph.text:
                    # Find the shape containing this text and update
                    text_frame = shape.text_frame
                    for i, p in enumerate(text_frame.paragraphs):
                        if "MFDS 의료기기 인증 (한국)" in p.text:
                            p.text = "MFDS Medical Device Approval (South Korea)"
                            p.font.size = Pt(13.5)
                        elif "HSA 승인 (싱가포르)" in p.text:
                            p.text = "HSA Approval (Singapore)"
                            p.font.size = Pt(13.5)

                    # Add Malaysia approval
                    new_p = text_frame.add_paragraph()
                    new_p.text = "MDA Approval (Malaysia)"
                    new_p.font.size = Pt(13.5)
                    new_p.level = 0

                # Update Research Output section with publication titles
                if "Gastrointestinal Endoscopy 등 국제 학술지 3편 게재" in paragraph.text:
                    paragraph.text = "5 publications in top-tier journals:"
                    paragraph.font.size = Pt(13.5)

                    # Add publication titles
                    text_frame = shape.text_frame

                    pub1 = text_frame.add_paragraph()
                    pub1.text = "• GI Endoscopy: Explainable CAD for Colorectal Polyps"
                    pub1.font.size = Pt(11)
                    pub1.level = 0

                    pub2 = text_frame.add_paragraph()
                    pub2.text = "• Clinical Endoscopy: AI-Assisted Colonoscopy for Adenoma"
                    pub2.font.size = Pt(11)
                    pub2.level = 0

                    pub3 = text_frame.add_paragraph()
                    pub3.text = "• Gastric Cancer: AI for Pathologic Outcome Prediction"
                    pub3.font.size = Pt(11)
                    pub3.level = 0

                    pub4 = text_frame.add_paragraph()
                    pub4.text = "• PMC: Real-World AI for Gastric Atypia Detection"
                    pub4.font.size = Pt(11)
                    pub4.level = 0

                    pub5 = text_frame.add_paragraph()
                    pub5.text = "• Cancer: AI Detection of Borrmann Type 4 Gastric Cancer"
                    pub5.font.size = Pt(11)
                    pub5.level = 0

    print("✓ Slide 6 updated with Malaysia approval and publication titles")

    # ============================================
    # Update Slide 7 (Phase 4: Impact) - Index 6
    # ============================================
    slide_7 = prs.slides[6]

    print("Updating Slide 7 (Phase 4: Impact)...")

    for shape in slide_7.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                # Update country count
                if "3개국" in paragraph.text:
                    paragraph.text = "4 Countries"
                    paragraph.font.bold = True
                elif "해외 진출" in paragraph.text:
                    paragraph.text = "International Expansion"
                # Update hospital count display
                elif "7개 병원" in paragraph.text:
                    paragraph.text = "7+ Hospitals"
                    paragraph.font.bold = True
                elif "국내외 배포" in paragraph.text:
                    paragraph.text = "Domestic & International"
                # Update Clinical Impact section
                elif "일평균 100건 이상 실시간 진단 지원" in paragraph.text:
                    paragraph.text = "100+ real-time diagnostic support cases daily"
                    paragraph.font.size = Pt(13.5)
                elif "용종 발견율 25% → 95% 이상 향상" in paragraph.text:
                    paragraph.text = "Polyp detection rate improved from 25% to 95%+"
                    paragraph.font.size = Pt(13.5)
                elif "의사 만족도 95%" in paragraph.text:
                    paragraph.text = "95% physician satisfaction rate"
                    paragraph.font.size = Pt(13.5)
                elif "Clinical Impact" in paragraph.text:
                    # Keep English title
                    pass

    print("✓ Slide 7 updated with 4 countries and real-time support")

    # ============================================
    # Update Slide 10 (Leadership & Scale) - Index 9
    # ============================================
    slide_10 = prs.slides[9]

    print("Updating Slide 10 (Leadership & Scale)...")

    for shape in slide_10.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                # Update Team Leadership section
                if "8명 개발팀 관리" in paragraph.text:
                    paragraph.text = "Led 8-member development team"
                    paragraph.font.size = Pt(13.5)
                elif "2년간 3개 제품 출시" in paragraph.text:
                    paragraph.text = "Launched 3 products in 2 years"
                    paragraph.font.size = Pt(13.5)
                elif "주니어 육성" in paragraph.text:
                    paragraph.text = "Junior developer mentorship"
                    paragraph.font.size = Pt(13.5)
                # Update Global Expansion section
                elif "싱가포르, 베트남, 태국" in paragraph.text:
                    paragraph.text = "Thailand, Singapore, Vietnam, Malaysia"
                    paragraph.font.size = Pt(13.5)
                elif "현지 병원 직접 방문 데모" in paragraph.text:
                    paragraph.text = "On-site hospital demonstrations and deployment"
                    paragraph.font.size = Pt(13.5)
                elif "기술지원" in paragraph.text:
                    paragraph.text = "Real-time diagnostic support and technical assistance"
                    paragraph.font.size = Pt(13.5)
                # Update Recognition section - keep as is
                elif "Team Leadership" in paragraph.text or "Global Expansion" in paragraph.text or "Recognition" in paragraph.text:
                    # Keep English section titles
                    pass

    print("✓ Slide 10 updated with 4 countries and enhanced descriptions")

    # Save the updated presentation
    output_path = r"D:\01_code\developer_portfolio\Jinbae_Park_Product_Leader.pptx"
    prs.save(output_path)

    print("\n[SUCCESS] All slides updated successfully!")
    print("[SUCCESS] Added Malaysia approval")
    print("[SUCCESS] Added publication titles")
    print("[SUCCESS] Updated to 4 countries")
    print("[SUCCESS] Added real-time diagnostic support information")
    print(f"[SUCCESS] Saved to: {output_path}")

except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
